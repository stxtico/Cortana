"""write_pptx (PROMPTS.md A24) - write-capable, gated the same way
tools/write_file.py already is: REQUIRES_CONFIRMATION = True, whitelisted to
[tools].write_whitelist_dirs via tools/_fs.py.
"""

from pptx import Presentation

from tools import _fs

REQUIRES_CONFIRMATION = True

_WRITE_KEY = "write_whitelist_dirs"
_TITLE_AND_CONTENT_LAYOUT = 1


def spec() -> dict:
    return {
        "type": "function",
        "function": {
            "name": "write_pptx",
            "description": (
                "Create a PowerPoint presentation (.pptx) from a list of slides, each with a "
                "title and body text (bullet lines separated by newlines). Only these "
                f"directories (and their subdirectories) are write-authorized: "
                f"{_fs.whitelist_description(_WRITE_KEY)}. Requires user confirmation before "
                "it actually runs."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to the .pptx file, absolute or relative to the project root."},
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "content": {"type": "string", "description": "Bullet lines, separated by newlines."},
                            },
                            "required": ["title"],
                        },
                        "description": "One entry per slide.",
                    },
                },
                "required": ["path", "slides"],
            },
        },
    }


def describe(path: str, slides: list) -> str:
    titles = ", ".join(s.get("title", "(untitled)") for s in slides[:5])
    more = f" +{len(slides) - 5} more" if len(slides) > 5 else ""
    return f"Write a {len(slides)}-slide presentation to {path!r}: {titles}{more}."


async def execute(path: str, slides: list) -> str:
    resolved = _fs.resolve_in_whitelist(path, _WRITE_KEY)
    resolved.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    layout = prs.slide_layouts[_TITLE_AND_CONTENT_LAYOUT]
    for slide_data in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "")
        content = slide_data.get("content") or ""
        if content and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            lines = content.split("\n")
            body.text_frame.text = lines[0]
            for line in lines[1:]:
                body.text_frame.add_paragraph().text = line
    prs.save(str(resolved))

    return f"Wrote a {len(slides)}-slide presentation to {resolved}."
